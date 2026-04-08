import os

import streamlit as st
from dotenv import load_dotenv

from architect import Provider, generate_agent_spec
from provider_vapi import deploy_vapi_assistant, get_phone_number
from supabase_client import create_deployed_agent, ensure_deployed_agents_table, list_deployed_agents
from runtime_agent import generate_chat_reply

load_dotenv()  # Load .env so OPENAI_API_KEY and others are available

st.set_page_config(page_title="AI Employee Studio", page_icon="🧱", layout="wide")

# Clean, modern, aurora-inspired light-ish UI
st.markdown(
    """
    <style>
    /* Global layout */
    .main {
        background: radial-gradient(circle at top, #ecfeff 0, #e0f2fe 35%, #f5f3ff 80%);
        color: #020617;
    }

    .block-container {
        padding-top: 3rem;
        padding-bottom: 3rem;
        max-width: 1200px;
    }

    /* Glass panels */
    .glass-card {
        background: linear-gradient(145deg, rgba(248,250,252,0.86), rgba(224,242,254,0.9));
        border-radius: 1.3rem;
        border: 1px solid rgba(148,163,184,0.35);
        box-shadow:
            0 18px 40px rgba(15,23,42,0.08),
            inset 0 0 0 1px rgba(255,255,255,0.55);
        backdrop-filter: blur(18px);
        -webkit-backdrop-filter: blur(18px);
        padding: 1.8rem 2rem;
    }

    .glass-pill {
        background: linear-gradient(120deg, rgba(45,212,191,0.4), rgba(59,130,246,0.45), rgba(236,72,153,0.45));
        border-radius: 999px;
        padding: 0.35rem 0.9rem;
        font-size: 0.8rem;
        letter-spacing: 0.08em;
        text-transform: uppercase;
        color: #0f172a;
        border: 1px solid rgba(148,163,184,0.45);
        display: inline-flex;
        align-items: center;
        gap: 0.4rem;
    }

    /* 3D-ish primary button */
    .stButton>button {
        width: 100%;
        border-radius: 999px;
        border: none;
        padding: 0.85rem 1.4rem;
        font-weight: 600;
        letter-spacing: 0.03em;
        text-transform: uppercase;
        color: #0f172a;
        background: linear-gradient(135deg, #a5b4fc, #67e8f9);
        box-shadow:
            0 14px 30px rgba(56,189,248,0.35),
            0 0 0 1px rgba(148,163,184,0.6);
        transform: translateY(0);
        transition: all 160ms ease-out;
    }
    .stButton>button:hover {
        filter: brightness(1.05);
        transform: translateY(-1px);
        box-shadow:
            0 18px 40px rgba(37,99,235,0.35),
            0 0 0 1px rgba(59,130,246,0.65);
    }
    .stButton>button:active {
        transform: translateY(1px);
        box-shadow:
            0 10px 20px rgba(15,23,42,0.25),
            0 0 0 1px rgba(148,163,184,0.8);
    }

    /* Text area + inputs */
    textarea, .stTextArea textarea {
        border-radius: 1rem !important;
        border: 1px solid rgba(148,163,184,0.6) !important;
        background: radial-gradient(circle at top left, rgba(255,255,255,0.95), rgba(224,242,254,0.96)) !important;
        color: #020617 !important;
        box-shadow: 0 12px 28px rgba(15,23,42,0.12);
    }

    /* Status component */
    [data-testid="stStatus"] {
        border-radius: 1rem;
        background: linear-gradient(135deg, rgba(239,246,255,0.98), rgba(224,242,254,0.96));
        border: 1px solid rgba(96,165,250,0.55);
        box-shadow: 0 18px 40px rgba(15,23,42,0.18);
    }

    /* Metric / info blocks */
    .metric-chip {
        background: radial-gradient(circle at top left, rgba(239,246,255,0.96), rgba(224,242,254,0.95));
        border-radius: 1rem;
        padding: 0.75rem 1rem;
        border: 1px solid rgba(148,163,184,0.55);
        font-size: 0.8rem;
        color: #0f172a;
    }
    .metric-chip span.label {
        text-transform: uppercase;
        letter-spacing: 0.12em;
        font-size: 0.7rem;
        color: #6b7280;
    }

    /* Hide default Streamlit footer */
    footer {visibility: hidden;}
    </style>
    """,
    unsafe_allow_html=True,
)

st.markdown(
    '<div class="glass-pill">AGENT BUILDER STUDIO</div>',
    unsafe_allow_html=True,
)
st.markdown(
    "<h1 style='margin-top: 0.8rem; font-size: 2.3rem;'>Design and deploy AI employees from a single brief.</h1>",
    unsafe_allow_html=True,
)
st.caption("Describe your business and ideal hire. Agent Builder Studio turns that into a voice-ready, tool-aware assistant.")

col_left, col_right = st.columns([2.1, 1.2], gap="large")

with col_left:
    with st.container():
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        st.markdown("#### Describe your business & vibe")
        description = st.text_area(
            " ",
            placeholder=(
                "e.g., I run a sushi restaurant in San Francisco and need a warm, patient receptionist who can "
                "answer menu questions, manage table bookings, and handle call overflow during dinner service."
            ),
            height=220,
            label_visibility="collapsed",
        )

        provider_label = st.selectbox(
            "LLM Provider",
            options=["OpenAI", "Claude"],
            index=0,
            help="Choose which model provider to use for architecting the agent spec.",
        )

        st.markdown("#### Agent capabilities")
        qualities = st.multiselect(
            "Qualities (tone & behavior)",
            options=[
                "friendly",
                "calm",
                "professional",
                "confident",
                "empathetic",
                "high-energy",
                "concise",
                "salesy (soft)",
                "detail-oriented",
                "strict about policy",
            ],
            default=["friendly", "professional", "concise"],
            help="Pick how the agent should sound and behave.",
        )

        abilities = st.multiselect(
            "Abilities (what it can do)",
            options=[
                "outbound_calls",
                "appointment_booking",
                "order_intake",
                "lead_qualification",
                "follow_up_reminders",
                "faq_support",
                "handoff_to_human",
                "collect_payments",
            ],
            default=["faq_support", "appointment_booking"],
            help="Choose functional abilities you want the agent to have.",
        )

        all_channels = st.checkbox(
            "Deploy on all channels (Phone + Telegram + WhatsApp + Instagram + LinkedIn + Web)",
            value=False,
        )
        channel_options = ["phone", "telegram", "whatsapp", "instagram", "linkedin", "web"]
        channels = (
            channel_options
            if all_channels
            else st.multiselect(
                "Channels",
                options=channel_options,
                default=["phone"],
                help="Choose where the agent should be available. (Voice deployment is live; others are captured as requirements.)",
            )
        )

        uploaded_files = st.file_uploader(
            "Upload reference files (optional)",
            type=["png", "jpg", "jpeg", "webp", "pdf", "ppt", "pptx", "docx"],
            accept_multiple_files=True,
            help="Upload menus, SOPs, brochures, contracts, policies, FAQs, etc. We'll extract text and use it as context.",
        )

        deploy_clicked = st.button("Deploy AI Employee")
        st.markdown("</div>", unsafe_allow_html=True)

with col_right:
    st.markdown('<div class="glass-card">', unsafe_allow_html=True)
    st.markdown("#### Deployment Snapshot")
    st.markdown(
        """
        <div class="metric-chip">
            <span class="label">Step 1</span><br/>
            Architect identity & tools from your description.
        </div>
        <br/>
        <div class="metric-chip">
            <span class="label">Step 2</span><br/>
            Deploy to Vapi.ai with a premium neural voice.
        </div>
        <br/>
        <div class="metric-chip">
            <span class="label">Step 3</span><br/>
            Persist assistant metadata to Supabase for analytics.
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.markdown("</div>", unsafe_allow_html=True)

st.markdown("---")

def _extract_text_from_uploads(files):
    extracted_chunks = []
    file_summaries = []
    images = []

    if not files:
        return "", file_summaries, images

    for f in files:
        name = getattr(f, "name", "upload")
        mime = getattr(f, "type", "") or ""
        data = f.getvalue()

        # Images: store bytes and show a short summary
        if mime.startswith("image/") or name.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
            images.append((name, data))
            file_summaries.append(f"{name} (image)")
            continue

        # PDFs
        if name.lower().endswith(".pdf"):
            try:
                from pypdf import PdfReader
                import io

                reader = PdfReader(io.BytesIO(data))
                text = "\n".join((page.extract_text() or "") for page in reader.pages)
                text = text.strip()
                file_summaries.append(f"{name} (pdf, {len(reader.pages)} pages)")
                if text:
                    extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n{text}")
                else:
                    extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n(No extractable text found)")
            except Exception as e:  # noqa: BLE001
                file_summaries.append(f"{name} (pdf, failed to parse)")
                extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n(Parse error: {e})")
            continue

        # DOCX
        if name.lower().endswith(".docx"):
            try:
                import io
                from docx import Document

                doc = Document(io.BytesIO(data))
                text = "\n".join(p.text for p in doc.paragraphs if p.text)
                file_summaries.append(f"{name} (docx)")
                extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n{text.strip()}")
            except Exception as e:  # noqa: BLE001
                file_summaries.append(f"{name} (docx, failed to parse)")
                extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n(Parse error: {e})")
            continue

        # PPTX (and PPT best-effort warning)
        if name.lower().endswith(".pptx"):
            try:
                import io
                from pptx import Presentation

                prs = Presentation(io.BytesIO(data))
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            t = (shape.text or "").strip()
                            if t:
                                texts.append(t)
                text = "\n".join(texts)
                file_summaries.append(f"{name} (pptx, {len(prs.slides)} slides)")
                extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n{text.strip()}")
            except Exception as e:  # noqa: BLE001
                file_summaries.append(f"{name} (pptx, failed to parse)")
                extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n(Parse error: {e})")
            continue

        if name.lower().endswith(".ppt"):
            file_summaries.append(f"{name} (ppt, not supported — please upload .pptx)")
            extracted_chunks.append(f"\n\n---\nSOURCE: {name}\n---\n(.ppt not supported; please export to .pptx)")
            continue

        file_summaries.append(f"{name} (unsupported)")

    return "\n".join(extracted_chunks).strip(), file_summaries, images


if deploy_clicked:
    if not description:
        st.error("Please provide a description first.")
    else:
        provider: Provider = "openai" if provider_label == "OpenAI" else "claude"
        extracted_text, file_summaries, images = _extract_text_from_uploads(uploaded_files)
        preferences = {
            "qualities": qualities,
            "abilities": abilities,
            "channels": channels,
        }

        with st.status("🧠 Architecting Agent Identity...", expanded=True) as status:
            try:
                spec = generate_agent_spec(
                    description,
                    provider=provider,
                    extracted_text=extracted_text,
                    file_summaries=file_summaries,
                    images=images,
                    preferences=preferences,
                )
                st.write(f"Created Identity: **{spec.name}**")

                status.update(label="📞 Provisioning Phone Line...", state="running")
                assistant_id = deploy_vapi_assistant(spec)
                phone_number = get_phone_number(assistant_id)
                st.session_state["assistant_id"] = assistant_id
                st.session_state["last_spec"] = spec
                print(f"[DEPLOY] assistant_id={assistant_id}")  # ensure visible in terminal logs

                # Persist deployment
                try:
                    ensure_deployed_agents_table()
                    create_deployed_agent(assistant_id=assistant_id, business_name=spec.name)
                except Exception as supa_err:  # noqa: BLE001
                    st.warning(f"Deployed, but DB insert failed: {supa_err}")

                status.update(label="🚀 Employee Deployed!", state="complete", expanded=False)

                st.success("Assistant deployed successfully")
                st.code(assistant_id)
                st.success(f"Call Now: {phone_number}")
                st.caption("Share this number with customers immediately.")

                with st.expander("Agent System Prompt", expanded=False):
                    st.code(spec.system_prompt, language="markdown")

                with st.expander("Tools Required", expanded=False):
                    st.json(spec.tools_required)

                with st.expander("Channels / Abilities / Qualities", expanded=False):
                    st.json(
                        {
                            "channels": getattr(spec, "channels", []),
                            "abilities": getattr(spec, "abilities", []),
                            "qualities": getattr(spec, "qualities", []),
                        }
                    )

            except Exception as e:  # noqa: BLE001
                st.error(f"Failed to deploy: {e}")

# Web channel: instant preview chat inside Streamlit (no extra setup)
st.markdown("### Web Chat Preview")
st.caption("Test your agent instantly in-browser. This powers the “web” channel behavior.")

if "preview_history" not in st.session_state:
    st.session_state["preview_history"] = []

preview_provider_label = st.selectbox(
    "Preview provider",
    options=["OpenAI", "Claude"],
    index=0,
    help="This controls which LLM answers in the web preview. It can differ from the architect provider.",
)
preview_provider: Provider = "openai" if preview_provider_label == "OpenAI" else "claude"

spec_for_preview = None
try:
    spec_for_preview = st.session_state.get("last_spec")
except Exception:
    spec_for_preview = None

if spec_for_preview is None:
    st.info("Deploy an agent above to enable the Web Chat Preview.")
else:
    for m in st.session_state["preview_history"]:
        with st.chat_message(m["role"]):
            st.markdown(m["content"])

    msg = st.chat_input("Type a customer message...")
    if msg:
        st.session_state["preview_history"].append({"role": "user", "content": msg})
        with st.chat_message("user"):
            st.markdown(msg)

        with st.chat_message("assistant"):
            reply = generate_chat_reply(
                spec=spec_for_preview,
                user_message=msg,
                history=st.session_state["preview_history"][:-1],
                provider=preview_provider,
            )
            st.markdown(reply)
        st.session_state["preview_history"].append({"role": "assistant", "content": reply})

st.markdown("### Channel Webhooks (Next)")
with st.expander("Enable Telegram / WhatsApp / Instagram / LinkedIn", expanded=False):
    st.markdown(
        """
Telegram webhook is available at:
- `/telegram/<assistant_id>`

Run webhook server:
```bash
python -m uvicorn server:app --host 0.0.0.0 --port 8000
```

Then expose it publicly and set Telegram webhook to:
- `https://<your-public-host>/telegram/<assistant_id>`

Also set:
- `TELEGRAM_BOT_TOKEN` in `.env`
        """
    )

st.markdown("### Deployed Agents")
st.caption("All deployed assistant IDs saved in Supabase.")
try:
    ensure_deployed_agents_table()
    deployed_agents_rows = list_deployed_agents(limit=200)
    if deployed_agents_rows:
        st.dataframe(deployed_agents_rows, use_container_width=True)
    else:
        st.info("No deployed agents found yet.")
except Exception as e:  # noqa: BLE001
    st.warning(f"Could not load deployed agents table: {e}")